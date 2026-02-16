"""
Document Intelligence - Core document processing engine.

Supports PDF, DOCX, PPTX, TXT with text extraction, chunking, and metadata.
"""

from typing import Dict, List, Any, Optional, BinaryIO
from dataclasses import dataclass
from pathlib import Path
import mimetypes
from datetime import datetime


@dataclass
class DocumentChunk:
    """A chunk of document content."""
    chunk_id: str
    content: str
    page_number: Optional[int] = None
    chunk_index: int = 0
    metadata: Dict[str, Any] = None
    

@dataclass
class Document:
    """Processed document with metadata."""
    doc_id: str
    filename: str
    file_type: str
    content: str
    chunks: List[DocumentChunk]
    metadata: Dict[str, Any]
    processed_at: str
    

class DocumentProcessor:
    """
    Core document processing engine for PDF, DOCX, PPTX, TXT files.
    
    Extracts text, creates chunks for RAG, and maintains metadata.
    """
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        """
        Initialize document processor.
        
        Args:
            chunk_size: Target size for each chunk (characters)
            chunk_overlap: Overlap between chunks for context preservation
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.supported_types = ['.pdf', '.docx', '.pptx', '.txt', '.md']
        
    def process_file(self, filepath: str) -> Document:
        """
        Process a document file.
        
        Args:
            filepath: Path to document file
            
        Returns:
            Processed Document object
        """
        path = Path(filepath)
        
        if not path.exists():
            raise FileNotFoundError(f"File not found: {filepath}")
        
        file_ext = path.suffix.lower()
        
        if file_ext not in self.supported_types:
            raise ValueError(f"Unsupported file type: {file_ext}. Supported: {self.supported_types}")
        
        if file_ext == '.pdf':
            content, metadata = self._process_pdf(filepath)
        elif file_ext == '.docx':
            content, metadata = self._process_docx(filepath)
        elif file_ext == '.pptx':
            content, metadata = self._process_pptx(filepath)
        elif file_ext in ['.txt', '.md']:
            content, metadata = self._process_text(filepath)
        else:
            raise ValueError(f"Handler not implemented for: {file_ext}")
        
        chunks = self._create_chunks(content)
        
        doc = Document(
            doc_id=self._generate_doc_id(filepath),
            filename=path.name,
            file_type=file_ext,
            content=content,
            chunks=chunks,
            metadata=metadata,
            processed_at=datetime.now().isoformat()
        )
        
        return doc
    
    def _process_pdf(self, filepath: str) -> tuple[str, Dict[str, Any]]:
        """Extract text from PDF."""
        try:
            import PyPDF2
            
            text_content = []
            metadata = {}
            
            with open(filepath, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                metadata = {
                    'pages': len(pdf_reader.pages),
                    'pdf_metadata': pdf_reader.metadata or {}
                }
                
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    text = page.extract_text()
                    if text.strip():
                        text_content.append(f"[Page {page_num}]\n{text}")
            
            return '\n\n'.join(text_content), metadata
            
        except ImportError:
            return self._process_pdf_fallback(filepath)
    
    def _process_pdf_fallback(self, filepath: str) -> tuple[str, Dict[str, Any]]:
        """Fallback PDF processing using pdfplumber."""
        try:
            import pdfplumber
            
            text_content = []
            metadata = {}
            
            with pdfplumber.open(filepath) as pdf:
                metadata = {
                    'pages': len(pdf.pages),
                    'pdf_metadata': pdf.metadata or {}
                }
                
                for page_num, page in enumerate(pdf.pages, 1):
                    text = page.extract_text()
                    if text and text.strip():
                        text_content.append(f"[Page {page_num}]\n{text}")
            
            return '\n\n'.join(text_content), metadata
            
        except ImportError:
            raise ImportError("PDF processing requires PyPDF2 or pdfplumber. Install: pip install PyPDF2 pdfplumber")
    
    def _process_docx(self, filepath: str) -> tuple[str, Dict[str, Any]]:
        """Extract text from DOCX."""
        try:
            from docx import Document as DocxDocument
            
            doc = DocxDocument(filepath)
            
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            
            tables_text = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join([cell.text for cell in row.cells])
                    if row_text.strip():
                        tables_text.append(row_text)
            
            content = '\n\n'.join(paragraphs)
            if tables_text:
                content += '\n\n[Tables]\n' + '\n'.join(tables_text)
            
            metadata = {
                'paragraphs': len(paragraphs),
                'tables': len(doc.tables),
                'sections': len(doc.sections)
            }
            
            return content, metadata
            
        except ImportError:
            raise ImportError("DOCX processing requires python-docx. Install: pip install python-docx")
    
    def _process_pptx(self, filepath: str) -> tuple[str, Dict[str, Any]]:
        """Extract text from PPTX."""
        try:
            from pptx import Presentation
            
            prs = Presentation(filepath)
            
            slides_text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = [f"[Slide {slide_num}]"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)
                
                if len(slide_content) > 1:
                    slides_text.append('\n'.join(slide_content))
            
            content = '\n\n'.join(slides_text)
            
            metadata = {
                'slides': len(prs.slides),
                'slide_width': prs.slide_width,
                'slide_height': prs.slide_height
            }
            
            return content, metadata
            
        except ImportError:
            raise ImportError("PPTX processing requires python-pptx. Install: pip install python-pptx")
    
    def _process_text(self, filepath: str) -> tuple[str, Dict[str, Any]]:
        """Process plain text files."""
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        metadata = {
            'lines': content.count('\n') + 1,
            'characters': len(content)
        }
        
        return content, metadata
    
    def _create_chunks(self, content: str) -> List[DocumentChunk]:
        """
        Split content into overlapping chunks for RAG.
        
        Args:
            content: Full document content
            
        Returns:
            List of DocumentChunk objects
        """
        if not content or not content.strip():
            return []
        
        chunks = []
        chunk_index = 0
        
        start = 0
        content_length = len(content)
        
        while start < content_length:
            end = start + self.chunk_size
            
            if end < content_length:
                end = self._find_sentence_boundary(content, end)
            else:
                end = content_length
            
            chunk_text = content[start:end].strip()
            
            if chunk_text:
                chunk = DocumentChunk(
                    chunk_id=f"chunk_{chunk_index}",
                    content=chunk_text,
                    chunk_index=chunk_index,
                    metadata={'start': start, 'end': end}
                )
                chunks.append(chunk)
                chunk_index += 1
            
            start = end - self.chunk_overlap
            
            if start >= content_length:
                break
        
        return chunks
    
    def _find_sentence_boundary(self, text: str, position: int) -> int:
        """Find nearest sentence boundary for clean chunk split."""
        search_window = 200
        search_start = max(0, position - search_window)
        search_text = text[search_start:position + search_window]
        
        sentence_endings = ['. ', '.\n', '! ', '!\n', '? ', '?\n']
        
        best_position = position
        
        for ending in sentence_endings:
            idx = search_text.rfind(ending, 0, position - search_start)
            if idx != -1:
                actual_pos = search_start + idx + len(ending)
                if abs(actual_pos - position) < abs(best_position - position):
                    best_position = actual_pos
        
        return best_position
    
    def _generate_doc_id(self, filepath: str) -> str:
        """Generate unique document ID."""
        import hashlib
        path_hash = hashlib.md5(filepath.encode()).hexdigest()[:12]
        return f"doc_{path_hash}"
    
    def get_chunk_by_id(self, document: Document, chunk_id: str) -> Optional[DocumentChunk]:
        """Retrieve specific chunk by ID."""
        for chunk in document.chunks:
            if chunk.chunk_id == chunk_id:
                return chunk
        return None
    
    def search_chunks(self, document: Document, query: str) -> List[DocumentChunk]:
        """Simple keyword search in chunks."""
        query_lower = query.lower()
        matching_chunks = []
        
        for chunk in document.chunks:
            if query_lower in chunk.content.lower():
                matching_chunks.append(chunk)
        
        return matching_chunks
